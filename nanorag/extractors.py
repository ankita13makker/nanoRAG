# Copyright (c) 2026, Salesforce, Inc.
# SPDX-License-Identifier: Apache-2.0

"""Text extraction from common file formats: PDF, DOCX, PPTX, XLSX, RTF, EPUB, etc."""

from __future__ import annotations

import re
from html.parser import HTMLParser
from pathlib import Path

TEXT_EXTENSIONS = {
    ".md", ".txt", ".py", ".java", ".js", ".ts", ".json", ".yaml", ".yml",
    ".xml", ".csv", ".html", ".htm", ".sql", ".sh", ".apex", ".cls",
    ".rst", ".log", ".ini", ".cfg", ".toml", ".r", ".rb", ".go",
    ".c", ".cpp", ".h", ".hpp", ".swift", ".kt", ".scala", ".pl",
}

BINARY_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".xlsx", ".rtf", ".epub", ".odt",
}

ALL_SUPPORTED = TEXT_EXTENSIONS | BINARY_EXTENSIONS


def is_supported(path: Path) -> bool:
    return path.suffix.lower() in ALL_SUPPORTED


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in TEXT_EXTENSIONS:
        return _read_text_file(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".rtf":
        return _extract_rtf(path)
    if ext == ".epub":
        return _extract_epub(path)
    if ext == ".odt":
        return _extract_odt(path)
    return ""


def _read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_pdf(path: Path) -> str:
    try:
        import pymupdf
    except ImportError:
        try:
            import fitz as pymupdf  # type: ignore[no-redef]
        except ImportError:
            raise ImportError(
                "PDF support requires PyMuPDF. Install it with: pip install pymupdf"
            )

    doc = pymupdf.open(str(path))
    pages = []
    for page in doc:
        text = page.get_text("text")
        if text.strip():
            pages.append(text.strip())
    doc.close()
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError:
        raise ImportError(
            "DOCX support requires python-docx. Install it with: pip install python-docx"
        )

    doc = docx.Document(str(path))
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name.lower() if para.style else ""
        if "heading" in style:
            level = 1
            for ch in style:
                if ch.isdigit():
                    level = int(ch)
                    break
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))

    return "\n\n".join(parts)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "PPTX support requires python-pptx. Install it with: pip install python-pptx"
        )

    prs = Presentation(str(path))
    slides: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        if texts:
            slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))

    return "\n\n".join(slides)


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError(
            "XLSX support requires openpyxl. Install it with: pip install openpyxl"
        )

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets: list[str] = []

    for ws in wb.worksheets:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"## {ws.title}\n\n" + "\n".join(rows))

    wb.close()
    return "\n\n".join(sheets)


def _extract_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ImportError(
            "RTF support requires striprtf. Install it with: pip install striprtf"
        )

    raw = path.read_bytes().decode("utf-8", errors="replace")
    return rtf_to_text(raw)


def _extract_epub(path: Path) -> str:
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        raise ImportError(
            "EPUB support requires ebooklib. Install it with: pip install ebooklib"
        )

    book = epub.read_epub(str(path), options={"ignore_ncx": True})
    parts: list[str] = []

    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        html = item.get_content().decode("utf-8", errors="replace")
        text = _strip_html(html)
        if text.strip():
            parts.append(text.strip())

    return "\n\n".join(parts)


def _extract_odt(path: Path) -> str:
    import zipfile
    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            content = zf.read("content.xml").decode("utf-8", errors="replace")
        return _strip_html(content)
    except Exception:
        return ""


class _HTMLStripper(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._skip = False

    def handle_starttag(self, tag: str, attrs: list) -> None:
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag: str) -> None:
        if tag in ("script", "style"):
            self._skip = False
        if tag in ("p", "div", "br", "h1", "h2", "h3", "h4", "h5", "h6", "li", "tr"):
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skip:
            self.parts.append(data)


def _strip_html(html: str) -> str:
    stripper = _HTMLStripper()
    stripper.feed(html)
    text = "".join(stripper.parts)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()
